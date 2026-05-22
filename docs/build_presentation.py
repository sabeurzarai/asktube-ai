"""Build a polished AskTube AI final project presentation.

The deck is generated from code so the final PPTX can be recreated after
content changes. The visual system is intentionally cinematic, dark, and
technical, matching the AskTube AI product direction.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PATH_OUT = ROOT / "docs" / "AskTube_AI_Presentation.pptx"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(5, 7, 13)
INK = RGBColor(246, 248, 255)
MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(34, 211, 238)
PINK = RGBColor(236, 72, 153)
PURPLE = RGBColor(139, 92, 246)
GREEN = RGBColor(52, 211, 153)
AMBER = RGBColor(251, 191, 36)
CARD = RGBColor(16, 23, 42)
CARD_2 = RGBColor(24, 31, 52)
LINE = RGBColor(51, 65, 85)


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_bg(slide, title: str | None = None) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Soft cinematic light fields.
    add_shape(slide, 0.0, 0.0, 13.333, 1.2, rgb("101827"), transparency=18, line=False)
    add_shape(slide, 8.4, -0.7, 5.2, 3.2, rgb("0e7490"), transparency=55, line=False)
    add_shape(slide, -1.0, 5.4, 4.0, 2.7, rgb("be185d"), transparency=62, line=False)

    if title:
        add_text(slide, title, 0.55, 0.35, 7.8, 0.35, size=11, color=MUTED, bold=True, caps=True, spacing=2.2)
        add_line(slide, 0.55, 0.82, 12.2, 0.82, LINE)


def add_shape(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    *,
    transparency: int = 0,
    line: bool = True,
    line_color: RGBColor = LINE,
    radius: bool = True,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line_color
        shape.line.transparency = 15
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = LINE, width: float = 1.2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 24,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
    caps: bool = False,
    spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text.upper() if caps else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if spacing is not None:
        run.font._element.set("spc", str(int(spacing * 100)))
    return box


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 15,
    color: RGBColor = INK,
    accent: RGBColor = CYAN,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.05
        r1 = p.add_run()
        r1.text = "- "
        r1.font.name = "Aptos"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = accent
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Aptos"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return box


def add_card(slide, label: str, body: str, x: float, y: float, w: float, h: float, accent: RGBColor = CYAN):
    add_shape(slide, x, y, w, h, CARD, transparency=4, line_color=accent)
    add_shape(slide, x, y, 0.08, h, accent, transparency=0, line=False, radius=False)
    add_text(slide, label, x + 0.22, y + 0.18, w - 0.35, 0.25, size=11, color=accent, bold=True, caps=True, spacing=1.6)
    add_text(slide, body, x + 0.22, y + 0.55, w - 0.35, h - 0.7, size=14, color=INK)


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, accent: RGBColor):
    add_text(slide, value, x, y, w, 0.65, size=36, color=accent, bold=True)
    add_text(slide, label, x, y + 0.62, w, 0.35, size=11, color=MUTED, bold=True, caps=True, spacing=1.4)


def add_stage(slide, number: str, title: str, body: str, x: float, y: float, accent: RGBColor):
    add_shape(slide, x, y, 2.25, 1.45, CARD_2, transparency=0, line_color=accent)
    add_shape(slide, x + 0.18, y + 0.18, 0.42, 0.42, accent, transparency=0, line=False)
    add_text(slide, number, x + 0.29, y + 0.27, 0.2, 0.2, size=12, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.72, y + 0.2, 1.35, 0.25, size=12, color=INK, bold=True)
    add_text(slide, body, x + 0.2, y + 0.72, 1.85, 0.45, size=10, color=MUTED)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, "AskTube AI", 0.7, 1.0, 7.6, 0.95, size=58, color=INK, bold=True)
    add_text(slide, "Chat with YouTube videos through transcript-grounded AI.", 0.75, 2.0, 7.8, 0.55, size=21, color=CYAN)
    add_text(slide, "Final Project | Sabeur Zarai | IronHack", 0.78, 6.78, 6.2, 0.28, size=12, color=MUTED, bold=True)

    add_shape(slide, 8.0, 1.0, 4.4, 4.8, CARD, transparency=0, line_color=CYAN)
    add_text(slide, "SEARCH", 8.35, 1.42, 3.7, 0.35, size=16, color=CYAN, bold=True, caps=True, spacing=2.5)
    add_text(slide, "TRANSCRIPT", 8.35, 2.16, 3.7, 0.35, size=16, color=PINK, bold=True, caps=True, spacing=2.5)
    add_text(slide, "RAG", 8.35, 2.9, 3.7, 0.45, size=25, color=INK, bold=True, caps=True, spacing=2.5)
    add_text(slide, "CITATIONS", 8.35, 3.72, 3.7, 0.35, size=16, color=GREEN, bold=True, caps=True, spacing=2.5)
    add_line(slide, 8.35, 4.45, 11.8, 4.45, PINK, 2)
    add_text(slide, "Netflix-inspired UI + FastAPI + LangChain + ChromaDB", 8.35, 4.72, 3.55, 0.5, size=13, color=MUTED)


def slide_problem_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "01 / Product Idea")
    add_text(slide, "YouTube is huge. Learning from long videos is slow.", 0.65, 1.15, 8.6, 0.75, size=34, color=INK, bold=True)
    add_text(slide, "AskTube AI turns a video into a searchable, conversational knowledge base.", 0.68, 2.0, 8.0, 0.45, size=18, color=CYAN)
    add_card(slide, "User flow", "Search a topic, choose a video, let the app process the transcript, then ask questions with source timestamps.", 0.72, 3.0, 3.55, 1.9, CYAN)
    add_card(slide, "Trust rule", "The assistant answers from transcript evidence only. If the transcript cannot answer, it refuses instead of inventing.", 4.55, 3.0, 3.55, 1.9, GREEN)
    add_card(slide, "Why it matters", "Learners can skip guessing where information lives and jump straight to the exact moment in the video.", 8.38, 3.0, 3.55, 1.9, PINK)


def slide_journey(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "02 / User Journey")
    add_text(slide, "From search to cited answer", 0.65, 1.0, 6.0, 0.6, size=34, color=INK, bold=True)
    stages = [
        ("1", "Search", "Text or voice search through YouTube Data API.", 0.7, 2.1, CYAN),
        ("2", "Choose", "Netflix-style carousel with selected video state.", 3.25, 2.1, PINK),
        ("3", "Process", "Transcript, chunks, embeddings, vector storage.", 5.8, 2.1, AMBER),
        ("4", "Chat", "RAG answer with citations and memory.", 8.35, 2.1, GREEN),
    ]
    for stage in stages:
        add_stage(slide, *stage)
    for x in [2.98, 5.53, 8.08]:
        add_line(slide, x, 2.82, x + 0.22, 2.82, CYAN, 2)
    add_text(slide, "The UI is intentionally cinematic: dark mode, glass surfaces, smooth transitions, a 3D assistant, and graceful loading/error states.", 0.8, 5.25, 11.4, 0.55, size=18, color=MUTED)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "03 / Architecture")
    add_text(slide, "Three layers, one learning loop", 0.65, 0.98, 7.4, 0.6, size=34, color=INK, bold=True)
    add_card(slide, "Frontend", "Next.js 14, TypeScript, TailwindCSS, shadcn/ui, Framer Motion, Embla, Three.js, React Three Fiber.", 0.75, 2.05, 3.45, 2.25, CYAN)
    add_card(slide, "Backend API", "FastAPI async routes for search, transcripts, chunks, vectorstore, chat, agent, speech, and evaluations.", 4.95, 2.05, 3.45, 2.25, PINK)
    add_card(slide, "AI + Data", "LangChain tools and agent, OpenAI models, ChromaDB vector storage, LangSmith tracing, Whisper fallback.", 9.15, 2.05, 3.45, 2.25, GREEN)
    add_line(slide, 4.2, 3.15, 4.9, 3.15, CYAN, 2)
    add_line(slide, 8.4, 3.15, 9.1, 3.15, PINK, 2)
    add_text(slide, "Docker Compose runs frontend, backend, and ChromaDB together for local development and deployment testing.", 1.0, 5.35, 10.8, 0.45, size=17, color=MUTED)


def slide_ingestion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "04 / Ingestion Pipeline")
    add_text(slide, "A YouTube video becomes retrieval-ready text.", 0.65, 0.98, 8.4, 0.6, size=33, color=INK, bold=True)
    items = [
        "YouTube Data API retrieves metadata, thumbnails, duration, and channel information.",
        "youtube-transcript-api is the primary source for public captions and timestamps.",
        "Whisper is used only as a fallback when captions are unavailable.",
        "LangChain chunking keeps timestamp metadata attached to every chunk.",
        "OpenAI embeddings are stored in ChromaDB for similarity search."
    ]
    add_bullets(slide, items, 0.85, 2.0, 6.6, 3.8, size=17, accent=AMBER)
    add_shape(slide, 8.15, 1.85, 3.8, 3.75, CARD, transparency=0, line_color=AMBER)
    add_text(slide, "Status stream", 8.48, 2.2, 3.0, 0.35, size=15, color=AMBER, bold=True, caps=True, spacing=1.5)
    add_text(slide, "metadata\ntranscript\nchunking\nembeddings\nvector_storage\nready / error", 8.5, 2.75, 2.9, 2.2, size=19, color=INK)
    add_text(slide, "The processing screen reflects backend progress, not only decoration.", 8.5, 5.1, 3.0, 0.36, size=11, color=MUTED)


def slide_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "05 / LangChain Agent")
    add_text(slide, "The chatbot uses tools, not just a single prompt.", 0.65, 0.98, 8.8, 0.6, size=33, color=INK, bold=True)
    add_text(slide, "AgentService binds StructuredTool objects to the LLM and returns answer, citations, tool steps, and session_id.", 0.72, 1.75, 10.8, 0.36, size=16, color=MUTED)

    tool_names = [
        "search_youtube_videos",
        "extract_transcript",
        "chunk_transcript",
        "store_video_vectors",
        "ingest_video",
        "retrieve_context",
        "answer_question",
    ]
    y = 2.35
    for idx, name in enumerate(tool_names):
        accent = [CYAN, PINK, AMBER, GREEN, PURPLE, CYAN, GREEN][idx]
        add_shape(slide, 0.85 + (idx % 2) * 3.55, y + (idx // 2) * 0.72, 3.15, 0.48, CARD_2, transparency=0, line_color=accent)
        add_text(slide, name, 1.05 + (idx % 2) * 3.55, y + 0.12 + (idx // 2) * 0.72, 2.75, 0.22, size=12, color=INK, bold=True)

    add_card(slide, "Guardrails", "The final answer must use answer_question, stay transcript-only, and include timestamp citations. Off-topic questions should be refused.", 8.15, 2.35, 3.85, 2.25, PINK)
    add_card(slide, "Memory", "ConversationMemoryService keeps session-based history so follow-up questions remain coherent.", 8.15, 4.85, 3.85, 1.15, CYAN)


def slide_rag(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "06 / RAG Answering")
    add_text(slide, "Answers are grounded in transcript evidence.", 0.65, 0.98, 8.3, 0.6, size=33, color=INK, bold=True)
    add_stage(slide, "1", "Retrieve", "ChromaDB returns top transcript chunks.", 0.85, 2.2, CYAN)
    add_stage(slide, "2", "Inject", "Prompt receives context, memory, and rules.", 3.55, 2.2, PINK)
    add_stage(slide, "3", "Generate", "GPT-4o-mini writes the answer.", 6.25, 2.2, AMBER)
    add_stage(slide, "4", "Cite", "Timestamp chips point to source moments.", 8.95, 2.2, GREEN)
    for x in [3.25, 5.95, 8.65]:
        add_line(slide, x, 2.92, x + 0.22, 2.92, CYAN, 2)
    add_text(slide, "If the transcript does not contain the answer, the system says it cannot answer from the video. This is the main anti-hallucination behavior.", 0.95, 5.25, 10.8, 0.5, size=18, color=MUTED)


def slide_voice_ux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "07 / Voice + Cinematic UX")
    add_text(slide, "The interface feels like a learning product, not a backend demo.", 0.65, 0.98, 9.5, 0.6, size=32, color=INK, bold=True)
    add_card(slide, "Voice search", "Browser Web Speech API captures live dictation with waveform feedback. If it fails, MediaRecorder sends audio to /api/speech/transcribe.", 0.75, 2.1, 3.65, 2.0, CYAN)
    add_card(slide, "Assistant", "A floating 3D robot moves through the journey, speaks short guidance, and opens a compact transcript-grounded chat.", 4.85, 2.1, 3.65, 2.0, PURPLE)
    add_card(slide, "Accessibility", "Semantic labels, visible focus states, reduced-motion handling, screen-reader status messages, and touch-friendly layouts.", 8.95, 2.1, 3.65, 2.0, GREEN)
    add_text(slide, "Extra polish: loading skeletons, retry states, TTS answer playback, responsive stacking, and collapsible evidence panels.", 0.95, 5.25, 10.8, 0.5, size=18, color=MUTED)


def slide_evaluation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "08 / Testing + Evaluation")
    add_text(slide, "The project is testable, explainable, and demo-ready.", 0.65, 0.98, 8.8, 0.6, size=33, color=INK, bold=True)
    add_metric(slide, "82", "pytest tests", 0.9, 2.15, 2.0, CYAN)
    add_metric(slide, "17", "RAG eval cases", 3.35, 2.15, 2.2, PINK)
    add_metric(slide, "0", "behavioral fails", 6.0, 2.15, 2.2, GREEN)
    add_metric(slide, "3", "Docker services", 8.45, 2.15, 2.2, AMBER)
    add_bullets(
        slide,
        [
            "Coverage includes routes, services, tools, WebSocket ingest, speech transcription, and RAG behavior.",
            "Evaluation cases test answerable questions, refusals, citation accuracy, summaries, memory, and hallucination prevention.",
            "LangSmith tracing can inspect latency, context, answer quality, and chain behavior."
        ],
        0.95,
        4.15,
        11.4,
        1.75,
        size=15,
        accent=GREEN,
    )


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "09 / Live Demo")
    add_text(slide, "Demo path: search, prepare, ask, verify.", 0.65, 0.98, 8.8, 0.6, size=33, color=INK, bold=True)
    add_bullets(
        slide,
        [
            "Open localhost:3000 and search: python tutorial for beginners.",
            "Pick a video and click Prepare to show real ingestion progress.",
            "Ask: What is Python used for?",
            "Show timestamp citations and the tool breadcrumb.",
            "Ask a follow-up to demonstrate memory.",
            "Ask an unrelated question to show refusal behavior.",
            "Use mic input and Read aloud if time allows."
        ],
        0.95,
        2.0,
        6.1,
        4.3,
        size=17,
        accent=CYAN,
    )
    add_shape(slide, 8.05, 2.05, 3.9, 3.75, CARD, transparency=0, line_color=PINK)
    add_text(slide, "What to emphasize", 8.38, 2.38, 3.2, 0.35, size=14, color=PINK, bold=True, caps=True, spacing=1.5)
    add_text(slide, "This is not a generic chatbot. It is a transcript-grounded learning workflow with tools, memory, vector search, citations, and a polished UI.", 8.38, 3.0, 3.1, 1.45, size=18, color=INK)
    add_text(slide, "Keep the demo focused on trust: every answer links back to the video.", 8.38, 4.9, 3.1, 0.55, size=12, color=MUTED)


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, "AskTube AI makes YouTube learnable by conversation.", 0.8, 1.2, 10.8, 1.25, size=42, color=INK, bold=True)
    add_text(slide, "Search videos. Build transcript knowledge. Ask questions. Verify with timestamps.", 0.85, 2.65, 9.7, 0.5, size=21, color=CYAN)
    add_card(slide, "Mandatory requirements", "LLM chatbot, tools, memory, vector database, text processing, UI, tests, evaluation, and Docker deployment are covered.", 0.85, 4.1, 3.8, 1.35, GREEN)
    add_card(slide, "Optional features", "Voice input, Whisper fallback, WebSocket streaming, TTS, 3D assistant, and LangSmith tracing are included.", 4.95, 4.1, 3.8, 1.35, PINK)
    add_card(slide, "Next improvements", "Persistent memory, streaming agent responses, user accounts, multi-video comparison, and hosted deployment.", 9.05, 4.1, 3.4, 1.35, CYAN)
    add_text(slide, "Thank you", 0.85, 6.72, 3.0, 0.3, size=14, color=MUTED, bold=True)


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_problem_solution(prs)
    slide_journey(prs)
    slide_architecture(prs)
    slide_ingestion(prs)
    slide_agent(prs)
    slide_rag(prs)
    slide_voice_ux(prs)
    slide_evaluation(prs)
    slide_demo(prs)
    slide_close(prs)

    prs.save(PATH_OUT)
    print(f"Saved: {PATH_OUT}")


if __name__ == "__main__":
    build()
