from pptx import Presentation
from pptx.util import Pt
from lxml import etree

NOTES = [
    # Slide 1 - Opening (30 sec)
    """[30 sec] AskTube AI is my final project.

The idea is simple: the user searches for a YouTube video, the app processes the transcript, and then the user can chat with the video content.

The answers are not generic AI answers — they are grounded in the transcript and include timestamp citations, so the user can verify exactly where the answer came from.

KEY MESSAGE: This is a learning assistant for YouTube videos, not just a normal chatbot.""",

    # Slide 2 - Product Idea (45 sec)
    """[45 sec] The problem I wanted to solve is that YouTube has a lot of good learning content, but long videos are difficult to search inside.

AskTube AI turns the video transcript into a searchable knowledge base. The user searches, chooses a video, the backend processes the transcript, and then the user can ask questions.

The trust rule is very important: if the transcript does not contain the answer, the assistant refuses instead of inventing.

KEY MESSAGE: The value is speed plus trust.""",

    # Slide 3 - User Journey (40 sec)
    """[40 sec] The user journey has four steps.

1. Search by text or voice.
2. Choose a video from a Netflix-style carousel.
3. The backend extracts the transcript, chunks it, embeds it, and stores it in ChromaDB.
4. The user chats with the video and gets timestamped answers.

I added duration filters, loading states, smooth animations, and a 3D assistant to make the flow feel like a real product.

KEY MESSAGE: The technical pipeline is hidden behind a simple user flow.""",

    # Slide 4 - Architecture (55 sec)
    """[55 sec] The architecture has three main layers.

Frontend: Next.js with TypeScript, Tailwind, Framer Motion, Embla carousel, and Three.js for the 3D assistant.

Backend: FastAPI with routes for search, transcripts, chunking, vector storage, chat, agent chat, speech, and evaluations.

AI and data layer: LangChain, OpenAI models, ChromaDB, Whisper fallback, and optional LangSmith tracing.

Everything runs with Docker Compose. There is also an observability layer: analytics events, RAG metrics, pipeline metrics, Prometheus metrics, and an /analytics dashboard.

KEY MESSAGE: Separated like a production app — UI, API, AI/data pipeline, and observability.""",

    # Slide 5 - Ingestion Pipeline (55 sec)
    """[55 sec] This is the most important backend flow.

When the user selects a video, the backend first gets metadata from the YouTube Data API. Then it tries to get the transcript using youtube-transcript-api — this was also the approach recommended by Carlos.

The transcript is split into chunks. Each chunk keeps timestamp metadata so the answer can cite exact moments later.

Then OpenAI embeddings are generated and stored in ChromaDB. After this step, the video is ready for RAG questions.

Whisper and yt-dlp exist only as a fallback if captions are unavailable. The normal flow does not download full videos.

KEY MESSAGE: Transcript first, audio fallback only when needed.""",

    # Slide 6 - LangChain Agent (55 sec)
    """[55 sec] To match the tool requirement, I implemented a LangChain tools layer.

Each tool wraps an existing backend service:
- search_youtube_videos
- extract_transcript
- chunk_transcript
- store_video_vectors
- retrieve_context
- answer_question

The AgentService binds these tools to the LLM. So the system is not just one prompt — it can decide which tool is needed and return the final answer with citations and the session ID.

There is also memory, so follow-up questions can use the previous conversation.

KEY MESSAGE: This directly covers the "LLM with tools and memory" requirement.""",

    # Slide 7 - RAG Answering (45 sec)
    """[45 sec] The RAG flow has four steps.

1. The user question retrieves the most relevant transcript chunks from ChromaDB.
2. Those chunks, the memory, and the rules are injected into the prompt.
3. The model writes the answer only from that context.
4. The UI shows timestamp citations.

If the answer is not in the transcript, the assistant says it cannot answer from the video. This is the main anti-hallucination rule.

KEY MESSAGE: Retrieval controls what the model is allowed to answer.""",

    # Slide 8 - Voice + Cinematic UX (35 sec)
    """[35 sec] I added optional features to make the project feel more complete.

Voice search uses the browser Web Speech API first. If that fails, the app records audio and sends it to Whisper for transcription.

There is also text-to-speech for AI answers, a 3D assistant scene, animated loading states, responsive layouts, and accessibility improvements — labels, focus states, and reduced-motion support.

KEY MESSAGE: Optional features support the learning experience, but the core is still RAG.""",

    # Slide 9 - Analytics + Observability (35 sec)
    """[35 sec] I added an analytics and observability system so the app is not only functional, but measurable.

Frontend tracks: searches, video selections, voice search, carousel use, suggested prompts, transcript opens, and timestamp clicks.

Backend tracks RAG metrics: retrieval latency, generation latency, chunks retrieved, token estimates, citation coverage, and hallucination warnings.

There is a dashboard at /analytics and Prometheus metrics at /metrics.

KEY MESSAGE: The project can explain what users do, how the RAG system behaves, and where the pipeline is slow.""",

    # Slide 10 - Testing + Evaluation (35 sec)
    """[35 sec] For testing, the backend has 98 pytest tests covering routes, services, tools, WebSocket ingestion, speech transcription, memory, vector storage, and RAG behavior.

I also created a RAG evaluation dataset with 17 cases testing answerable questions, refusals, citation accuracy, summaries, memory, and hallucination prevention.

LangSmith tracing can also be enabled to inspect latency, context, tool calls, and answer quality.

KEY MESSAGE: I did not only build the app — I also tested the AI behavior.""",

    # Slide 11 - Live Demo (40 sec)
    """[40 sec] For the demo I will show the main flow: search, prepare, ask, and verify.

Steps:
1. Search "python tutorial for beginners" with a duration filter.
2. Pick a video and click Prepare — show the real ingestion progress.
3. Ask: "What is Python used for?" — show timestamp citations.
4. Ask a follow-up to demonstrate memory.
5. Ask an unrelated question to show refusal behavior.
6. Show /analytics and /metrics if time allows.

I can show the EC2 deployment to prove Docker works on a server. For the full transcript/RAG demo, local Docker is more reliable because YouTube blocks transcript requests from cloud IP ranges.

KEY MESSAGE: EC2 proves deployment; local demo proves full transcript ingestion.""",

    # Slide 12 - Closing (30 sec)
    """[30 sec] To summarize, AskTube AI meets the main project requirements: LLM chatbot, tools, memory, vector database, text processing, UI, tests, evaluation, and deployment.

It also includes optional features: voice input, Whisper fallback, streaming, TTS, 3D assistant, analytics, Prometheus metrics, and LangSmith tracing.

The main future improvements: persistent user accounts, multi-video comparison, a paid custom domain, and a stronger production transcript proxy.

FINAL LINE: AskTube AI makes YouTube videos learnable by conversation, while keeping the answer grounded in the original transcript.""",
]


def set_notes(slide, text):
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
    else:
        notes_slide = slide.notes_slide  # auto-creates it

    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.text = text


def main():
    path = "AskTube_AI_Presentation_updated.pptx"
    prs = Presentation(path)

    for i, (slide, note) in enumerate(zip(prs.slides, NOTES)):
        set_notes(slide, note)
        print(f"Added notes to slide {i + 1}")

    prs.save(path)
    print(f"\nSaved: {path}")


if __name__ == "__main__":
    main()
